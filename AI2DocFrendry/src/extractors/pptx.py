"""PowerPoint (.pptx) 抽出器。

スライドはシェイプの配置順＝論理順とは限らない。XML 順のまま読むと
文が入れ替わり、下流には「順序だけ壊れた自然な文章」として届いて検知できない。
そのため本文は (top, left) で並べ替えてから抽出する（FR-09 / ADR-03）。
"""

from __future__ import annotations

from typing import Any, Iterator

from ..models import Block, ImageRef, RawDocument, SourceFile, StyleMark
from .base import Extractor, ExtractionError, build, expand_merged

_IMPORT_ERROR = ""
try:
    from pptx import Presentation as _open_presentation
except Exception as _exc:  # ライブラリ未導入でも import 時に落とさない（NFR-13）
    _open_presentation = None  # type: ignore[assignment]
    _IMPORT_ERROR = str(_exc)

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_BULLET_TAGS = (f"{{{_A}}}buChar", f"{{{_A}}}buAutoNum")
_PPR_TAG = f"{{{_A}}}pPr"
_PICTURE_TYPES = ("PICTURE", "LINKED_PICTURE", "MEDIA")
_FIGURE_TYPES = ("AUTO_SHAPE", "FREEFORM", "DIAGRAM", "CANVAS", "LINE")


class PptxExtractor:
    """python-pptx を用いた .pptx 抽出器。"""

    extensions: tuple[str, ...] = (".pptx",)

    def extract(self, source: SourceFile) -> RawDocument:
        """.pptx を RawDocument に変換する。失敗は ExtractionError で送出する。"""
        if _open_presentation is None:
            raise ExtractionError(f"python-pptx が利用できません: {_IMPORT_ERROR}")
        blocks: list[Block] = []
        images: list[ImageRef] = []
        marks: list[StyleMark] = []
        try:
            presentation = _open_presentation(source.path)
            for number, slide in enumerate(presentation.slides, start=1):
                _read_slide(slide, f"slide{number}", blocks, images, marks)
        except ExtractionError:
            raise
        except Exception as exc:
            raise ExtractionError(f"pptx の抽出に失敗しました: {source.path}: {exc}") from exc
        return build(source, blocks, images, marks, _warnings(marks, images))


def _read_slide(
    slide: Any,
    position: str,
    blocks: list[Block],
    images: list[ImageRef],
    marks: list[StyleMark],
) -> None:
    """1スライドを読む。タイトルを先頭に固定し、本文は読み順に並べ替える。"""
    title = _title_shape(slide)
    title_element = getattr(title, "_element", None) if title is not None else None
    if title is not None:
        text = (title.text_frame.text or "").strip()
        if text:
            blocks.append(Block(kind="heading", text=text, level=1, position=position))
    body = [s for s in _iter_shapes(slide.shapes) if not _is_element(s, title_element)]
    for shape in sorted(body, key=_reading_key):
        _read_shape(shape, position, blocks, images, marks)
    note = _notes_text(slide)
    if note:
        blocks.append(Block(kind="footnote", text=note, position=position))


def _read_shape(
    shape: Any,
    position: str,
    blocks: list[Block],
    images: list[ImageRef],
    marks: list[StyleMark],
) -> None:
    """シェイプ1つを対応する要素へ振り分ける。"""
    if getattr(shape, "has_chart", False):
        images.append(ImageRef(position=position, kind="chart"))
        return
    if getattr(shape, "has_table", False):
        blocks.append(_table_block(shape.table, position))
        return
    if _is_visual(shape):
        images.append(ImageRef(position=position, kind="image"))
        return
    if not getattr(shape, "has_text_frame", False):
        return
    text_blocks = _text_blocks(shape.text_frame, position)
    if not text_blocks:
        return
    blocks.extend(text_blocks)
    if _is_textbox(shape):
        marks.append(_textbox_mark(text_blocks, position))


def _text_blocks(frame: Any, position: str) -> list[Block]:
    """テキストフレームの段落を Block 化する。空段落は捨てる。"""
    blocks: list[Block] = []
    for paragraph in frame.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        level = _para_level(paragraph)
        if level > 0 or _has_bullet(paragraph):
            blocks.append(Block(kind="list", text=text, level=level + 1, position=position))
        else:
            blocks.append(Block(kind="paragraph", text=text, position=position))
    return blocks


def _table_block(table: Any, position: str) -> Block:
    """表を Block にする。結合セルは矩形展開し行列対応を保つ（FR-07）。"""
    rows: list[list[str]] = []
    for row in table.rows:
        values: list[str] = []
        for index, cell in enumerate(row.cells):
            values.append(_cell_text(cell, rows, index))
        rows.append(values)
    return Block(kind="table", table_rows=expand_merged(rows), position=position)


def _cell_text(cell: Any, previous_rows: list[list[str]], index: int) -> str:
    """セル文字列。縦結合されたセルは空になるため直上の値を引き継ぐ。"""
    text = (cell.text or "").strip()
    if text or not getattr(cell, "is_spanned", False) or not previous_rows:
        return text
    above = previous_rows[-1]
    return above[index] if index < len(above) else ""


def _textbox_mark(text_blocks: list[Block], position: str) -> StyleMark:
    """テキストボックス由来である旨を残す。読み順の確からしさが本文と異なる。"""
    head = text_blocks[0].text
    return StyleMark(kind="textbox", text=head, position=position, preserved=True)


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """グループを再帰的に平坦化する。"""
    for shape in shapes:
        if _shape_type_name(shape) == "GROUP":
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _reading_key(shape: Any) -> tuple[int, int]:
    """読み順の推定キー。座標未設定のシェイプは末尾へ送る。"""
    top = getattr(shape, "top", None)
    left = getattr(shape, "left", None)
    return (
        int(top) if top is not None else 10**12,
        int(left) if left is not None else 10**12,
    )


def _title_shape(slide: Any) -> Any:
    """タイトルプレースホルダを返す。"""
    try:
        title = slide.shapes.title
    except Exception:
        return None
    if title is None or not getattr(title, "has_text_frame", False):
        return None
    return title


def _is_element(shape: Any, element: Any) -> bool:
    """同一シェイプ判定。python-pptx はアクセスの度に別プロキシを返すため XML 要素で比べる。"""
    if element is None:
        return False
    return getattr(shape, "_element", None) is element


def _notes_text(slide: Any) -> str:
    """ノートを取得する。前提や補足が本文外に置かれている場合がある。"""
    try:
        if not slide.has_notes_slide:
            return ""
        frame = slide.notes_slide.notes_text_frame
    except Exception:
        return ""
    if frame is None:
        return ""
    return (frame.text or "").strip()


def _para_level(paragraph: Any) -> int:
    """箇条書きの階層を返す。取得できない場合は 0 とみなす。"""
    fmt = getattr(paragraph, "paragraph_format", None)
    level = getattr(fmt, "level", None)
    if level is None:
        level = getattr(paragraph, "level", 0)
    try:
        return max(int(level), 0)
    except (TypeError, ValueError):
        return 0


def _has_bullet(paragraph: Any) -> bool:
    """行頭記号の有無を XML から判定する。level 0 でも箇条書きはあり得る。"""
    try:
        properties = paragraph._p.find(_PPR_TAG)
    except Exception:
        return False
    if properties is None:
        return False
    return any(properties.find(tag) is not None for tag in _BULLET_TAGS)


def _shape_type_name(shape: Any) -> str:
    try:
        shape_type = shape.shape_type
    except Exception:
        return ""
    return getattr(shape_type, "name", "") or ""


def _is_textbox(shape: Any) -> bool:
    if _shape_type_name(shape) == "TEXT_BOX":
        return True
    return not bool(getattr(shape, "is_placeholder", False))


def _is_visual(shape: Any) -> bool:
    """テキストを持たない図形・画像かどうか。"""
    name = _shape_type_name(shape)
    if name in _PICTURE_TYPES:
        return True
    if name not in _FIGURE_TYPES:
        return False
    if not getattr(shape, "has_text_frame", False):
        return True
    return not (shape.text_frame.text or "").strip()


def _warnings(marks: list[StyleMark], images: list[ImageRef]) -> list[str]:
    """読み順とチャートの不確実性を明示する。"""
    warnings: list[str] = []
    textboxes = sum(1 for m in marks if m.kind == "textbox")
    if textboxes:
        warnings.append(f"テキストボックス {textboxes} 個は読み順を座標から推定しています")
    charts = sum(ref.count for ref in images if ref.kind == "chart")
    if charts:
        warnings.append(f"グラフ {charts} 個は本文化されていません（OCR でも復元不可）")
    return warnings


#: 契約充足を型検査で担保する（実行時の副作用はない）
_CONTRACT: Extractor = PptxExtractor()
